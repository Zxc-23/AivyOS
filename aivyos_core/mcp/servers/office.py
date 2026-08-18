"""MCP office Server（文档 §5.1.2 / T3.6）：Word/Excel/PPT 生成。

- 真实后端：python-docx / openpyxl / python-pptx（可选）
- 零依赖回退：stdlib zipfile 生成**最小合法** docx/xlsx/pptx（OOXML 结构）
"""

from __future__ import annotations

import zipfile
from pathlib import Path
from typing import Any, Dict, Optional

from aivyos_core.mcp.types import PermissionLevel, Tool, ToolResult, make_tool

_CT = """<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<Types xmlns="http://schemas.openxmlformats.org/package/2006/content-types">\
<Default Extension="rels" ContentType="application/vnd.openxmlformats-package.relationships+xml"/>\
<Default Extension="xml" ContentType="application/xml"/>\
{extra}</Types>"""

_DOCX_RELS = """<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">\
<Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/officeDocument" Target="word/document.xml"/>\
</Relationships>"""

_DOCX_DOC = """<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<w:document xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main">\
<w:body><w:p><w:r><w:t xml:space="preserve">{text}</w:t></w:r></w:p></w:body></w:document>"""

_XLSX_SHEET = """<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<worksheet xmlns="http://schemas.openxmlformats.org/spreadsheetml/2006/main"><sheetData>{rows}</sheetData></worksheet>"""

_XLSX_WB = """<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<workbook xmlns="http://schemas.openxmlformats.org/spreadsheetml/2006/main" \
xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships">\
<sheets><sheet name="Sheet1" sheetId="1" r:id="rId1"/></sheets></workbook>"""


def _minimal_docx(text: str) -> bytes:
    buf = _zip({
        "[Content_Types].xml": _CT.format(extra='<Override PartName="/word/document.xml" ContentType="application/vnd.openxmlformats-officedocument.wordprocessingml.document.main+xml"/>'),
        "_rels/.rels": _DOCX_RELS,
        "word/document.xml": _DOCX_DOC.format(text=text[:1000]),
    })
    return buf


def _minimal_xlsx(cells: Dict[str, str]) -> bytes:
    rows_xml = ""
    for ref, value in cells.items():
        col = "".join(ch for ch in ref if ch.isalpha())
        row = "".join(ch for ch in ref if ch.isdigit())
        rows_xml += f'<row r="{row}"><c r="{ref}" t="inlineStr"><is><t>{value}</t></is></c></row>'
    sheet = _XLSX_SHEET.format(rows=rows_xml)
    return _zip({
        "[Content_Types].xml": _CT.format(extra=(
            '<Override PartName="/xl/workbook.xml" ContentType="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet.main+xml"/>'
            '<Override PartName="/xl/worksheets/sheet1.xml" ContentType="application/vnd.openxmlformats-officedocument.spreadsheetml.worksheet+xml"/>')),
        "_rels/.rels": '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
                       '<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">'
                       '<Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/officeDocument" Target="xl/workbook.xml"/></Relationships>',
        "xl/workbook.xml": _XLSX_WB,
        "xl/_rels/workbook.xml.rels": '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
                                      '<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">'
                                      '<Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/worksheet" Target="worksheets/sheet1.xml"/></Relationships>',
        "xl/worksheets/sheet1.xml": sheet,
    })


def _minimal_pptx(text: str) -> bytes:
    slide = ('<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
             '<p:sld xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main">'
             '<p:cSld><p:spTree><p:nvGrpSpPr><p:cNvPr id="1" name=""/><p:cNvGrpSpPr/><p:nvPr/></p:nvGrpSpPr>'
             '<p:grpSpPr/><p:sp><p:nvSpPr><p:cNvPr id="2" name="t"/><p:cNvSpPr/><p:nvPr/></p:nvSpPr>'
             '<p:spPr/><p:txBody><a:bodyPr xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"/>'
             '<a:lstStyle/><a:p xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
             f'<a:r><a:rPr lang="zh-CN"/><a:t>{text[:500]}</a:t></a:r></a:p></p:txBody></p:sp>'
             '</p:spTree></p:cSld></p:sld>')
    return _zip({
        "[Content_Types].xml": _CT.format(extra=(
            '<Override PartName="/ppt/presentation.xml" ContentType="application/vnd.openxmlformats-officedocument.presentationml.presentation.main+xml"/>'
            '<Override PartName="/ppt/slides/slide1.xml" ContentType="application/vnd.openxmlformats-officedocument.presentationml.slide+xml"/>')),
        "_rels/.rels": '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
                       '<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">'
                       '<Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/officeDocument" Target="ppt/presentation.xml"/></Relationships>',
        "ppt/presentation.xml": '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
                                '<p:presentation xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"><p:sldIdLst><p:sldId id="256" r:id="rId1"/></p:sldIdLst></p:presentation>',
        "ppt/_rels/presentation.xml.rels": '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
                                           '<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">'
                                           '<Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/slide" Target="slides/slide1.xml"/></Relationships>',
        "ppt/slides/slide1.xml": slide,
    })


def _zip(files: Dict[str, str]) -> bytes:
    import io

    buf = io.BytesIO()
    with zipfile.ZipFile(buf, "w", zipfile.ZIP_DEFLATED) as z:
        for name, content in files.items():
            z.writestr(name, content)
    return buf.getvalue()


class OfficeServer:
    def __init__(self, output_dir: Path) -> None:
        self.output = Path(output_dir)
        self.output.mkdir(parents=True, exist_ok=True)

    def _save(self, name: str, data: bytes) -> Path:
        path = self.output / name
        path.write_bytes(data)
        return path

    async def _create_docx(self, args: Dict[str, Any]) -> ToolResult:
        name = args.get("name", "文档.docx")
        text = args.get("text", "")
        try:
            from docx import Document  # type: ignore

            doc = Document()
            doc.add_paragraph(text)
            doc.save(self.output / name)
            return ToolResult(True, content=f"已生成: {name}（python-docx）", data={"path": str(self.output / name)})
        except ImportError:
            path = self._save(name, _minimal_docx(text))
            return ToolResult(True, content=f"已生成: {name}（stdlib 最小 docx）", data={"path": str(path)})

    async def _create_xlsx(self, args: Dict[str, Any]) -> ToolResult:
        name = args.get("name", "表格.xlsx")
        cells = args.get("cells") or {"A1": "值"}
        try:
            from openpyxl import Workbook  # type: ignore

            wb = Workbook()
            ws = wb.active
            for ref, value in cells.items():
                ws[ref] = value
            wb.save(self.output / name)
            return ToolResult(True, content=f"已生成: {name}（openpyxl）", data={"path": str(self.output / name)})
        except ImportError:
            path = self._save(name, _minimal_xlsx(cells))
            return ToolResult(True, content=f"已生成: {name}（stdlib 最小 xlsx）", data={"path": str(path)})

    async def _create_pptx(self, args: Dict[str, Any]) -> ToolResult:
        name = args.get("name", "演示.pptx")
        text = args.get("text", "")
        try:
            from pptx import Presentation  # type: ignore

            prs = Presentation()
            prs.slides.add_slide(prs.slide_layouts[1])
            prs.save(self.output / name)
            return ToolResult(True, content=f"已生成: {name}（python-pptx）", data={"path": str(self.output / name)})
        except ImportError:
            path = self._save(name, _minimal_pptx(text))
            return ToolResult(True, content=f"已生成: {name}（stdlib 最小 pptx）", data={"path": str(path)})

    def tools(self) -> list[Tool]:
        return [
            make_tool(
                "office_docx", "生成 Word 文档（L1）",
                {"type": "object", "properties": {"name": {"type": "string"}, "text": {"type": "string"}}, "required": ["text"]},
                self._create_docx, PermissionLevel.L1, server="office",
            ),
            make_tool(
                "office_xlsx", "生成 Excel 表格（L1）",
                {"type": "object", "properties": {"name": {"type": "string"}, "cells": {"type": "object"}}},
                self._create_xlsx, PermissionLevel.L1, server="office",
            ),
            make_tool(
                "office_pptx", "生成 PPT（L1）",
                {"type": "object", "properties": {"name": {"type": "string"}, "text": {"type": "string"}}},
                self._create_pptx, PermissionLevel.L1, server="office",
            ),
        ]
